import os
import sys
import pandas as pd
import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation
from typing import List
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_ollama import OllamaEmbeddings
from langchain_chroma import Chroma

# --- Path Optimization ---
_SERVICES_DIR = os.path.dirname(os.path.abspath(__file__))
_APP_DIR = os.path.dirname(_SERVICES_DIR)
_ROOT_DIR = os.path.dirname(_APP_DIR)

if _ROOT_DIR not in sys.path:
    sys.path.insert(0, _ROOT_DIR)
if _APP_DIR not in sys.path:
    sys.path.insert(0, _APP_DIR)

from services.config import MODELS_DIR, DOCS_DIR

# ---------------- CONFIG ---------------- #
PERSIST_DIR = os.path.join(MODELS_DIR, "files_vector")

SUPPORTED_TEXT_EXT = {
    ".txt", ".py", ".js", ".ts", ".jsx", ".tsx", ".html", ".css",
    ".json", ".yaml", ".yml", ".md", ".env", ".sql", ".sh",
    ".java", ".c", ".cpp", ".h", ".go", ".rb", ".php",
    ".xml", ".ini", ".conf", ".toml", ".bat", ".ps1"
}

SUPPORTED_OFFICE_EXT = {".pdf", ".docx", ".pptx", ".xlsx", ".xls"}
SUPPORTED_DATA_EXT = {".csv"}
ALL_SUPPORTED_EXT = SUPPORTED_TEXT_EXT | SUPPORTED_OFFICE_EXT | SUPPORTED_DATA_EXT

# ---------------- VECTOR STORE ---------------- #

embeddings = OllamaEmbeddings(model="nomic-embed-text:latest")

vector_store = Chroma(
    collection_name="smart_files",
    embedding_function=embeddings,
    persist_directory=PERSIST_DIR,
)

text_splitter = RecursiveCharacterTextSplitter(
    chunk_size=800,       
    chunk_overlap=150
)

# ---------------- TEXT EXTRACTION ---------------- #

def extract_text(filepath: str) -> List[Document]:
    ext = os.path.splitext(filepath)[1].lower()
    docs = []

    try:
        # -------- PDF -------- #
        if ext == ".pdf":
            try:
                with pdfplumber.open(filepath) as pdf:
                    for i, page in enumerate(pdf.pages):
                        text = page.extract_text() or ""
                        if text.strip():
                            docs.append(Document(
                                page_content=text,
                                metadata={"file_path": filepath, "page": i, "type": "pdf"}
                            ))
            except Exception as e:
                print(f"❌ Error processing {filepath}: {e}")

        # -------- DOCX -------- #
        elif ext == ".docx":
            doc = DocxDocument(filepath)
            text = "\n".join([p.text for p in doc.paragraphs])

            if text.strip():
                docs.append(Document(
                    page_content=text,
                    metadata={"file_path": filepath, "type": "docx"}
                ))

        # -------- TEXT -------- #
        elif ext in SUPPORTED_TEXT_EXT:
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()

            if text.strip():
                docs.append(Document(
                    page_content=text,
                    metadata={"file_path": filepath, "type": "text"}
                ))

        # -------- CSV -------- #
        elif ext == ".csv":
            df = pd.read_csv(filepath).head(5)
            text = df.to_string(index=False)

            docs.append(Document(
                page_content=text,
                metadata={"file_path": filepath, "type": "csv", "rows": len(df)}
            ))

        # -------- EXCEL -------- #
        elif ext in {".xlsx", ".xls"}:
            df = pd.read_excel(filepath)
            text = df.to_string(index=False)

            docs.append(Document(
                page_content=text,
                metadata={"file_path": filepath, "type": "excel", "rows": len(df)}
            ))

        # -------- PPTX -------- #
        elif ext == ".pptx":
            prs = Presentation(filepath)
            text_runs = []

            for i, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(f"[Slide {i+1}] {shape.text}")

            if text_runs:
                docs.append(Document(
                    page_content="\n".join(text_runs),
                    metadata={"file_path": filepath, "type": "pptx"}
                ))

    except Exception as e:
        print(f"❌ Error processing {filepath}: {e}")

    return docs


# ---------------- EMBEDDING ---------------- #

def embed_file(filepath: str):
    print(f"\nProcessing: {filepath}")

    docs = extract_text(filepath)

    if not docs:
        print("No content found, skipping...")
        return

    splits = text_splitter.split_documents(docs)

    splits = [s for s in splits if s.page_content.strip()]

    if not splits:
        print("No valid chunks, skipping...")
        return

    print(f"Chunks: {len(splits)}")

    batch_size = 50
    for i in range(0, len(splits), batch_size):
        vector_store.add_documents(splits[i:i+batch_size])

    print("Added to vector DB")

    return 1


# ---------------- FOLDER INGESTION ---------------- #

def embed_folder(folder_path: str):
    print(f"\n📁 Processing folder: {folder_path}")

    for root, _, files in os.walk(folder_path):
        for file in files:
            filepath = os.path.join(root, file)

            ext = os.path.splitext(filepath)[1].lower()
            if ext in ALL_SUPPORTED_EXT:
                embed_file(filepath)


if __name__ == "__main__":
    embed_folder(DOCS_DIR)